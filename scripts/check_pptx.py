from pptx import Presentation
prs = Presentation("outputs/HYMIND_Executive_Presentation.pptx")
print(f"Slides: {len(prs.slides)}")
print(f"Width:  {prs.slide_width.inches:.2f} in")
print(f"Height: {prs.slide_height.inches:.2f} in")
for i, sl in enumerate(prs.slides, 1):
    texts = [s.text.strip()[:60] for s in sl.shapes if hasattr(s, "text") and s.text.strip()]
    label = texts[0] if texts else "(no text)"
    print(f"  Slide {i:02d}: {label}")
