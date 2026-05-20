"""
HYMIND Executive Presentation Builder
Generates a professional 16:9 PowerPoint presentation.

Run:
    python scripts/build_presentation.py
Output:
    outputs/HYMIND_Executive_Presentation.pptx
"""

from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu, Cm
import pptx.util as util

# ---------------------------------------------------------------------------
# DESIGN SYSTEM
# ---------------------------------------------------------------------------

# Slide canvas — 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

# Color palette
C = {
    "bg_deep":      RGBColor(0x04, 0x0C, 0x1A),   # ultra-deep navy  #040C1A
    "bg_card":      RGBColor(0x09, 0x14, 0x2B),   # dark card        #09142B
    "bg_mid":       RGBColor(0x0E, 0x1E, 0x3E),   # mid navy         #0E1E3E
    "bg_section":   RGBColor(0x06, 0x10, 0x24),   # section bg       #061024
    "bg_light":     RGBColor(0xF4, 0xF7, 0xFC),   # off-white        #F4F7FC
    "cyan":         RGBColor(0x00, 0xD4, 0xFF),   # primary cyan     #00D4FF
    "cyan_dim":     RGBColor(0x00, 0x8C, 0xBD),   # dim cyan         #008CBD
    "blue":         RGBColor(0x10, 0x6B, 0xE8),   # electric blue    #106BE8
    "blue_dim":     RGBColor(0x0A, 0x42, 0x90),   # dim blue         #0A4290
    "white":        RGBColor(0xFF, 0xFF, 0xFF),
    "white_80":     RGBColor(0xCC, 0xD4, 0xE0),   # slightly dimmed  #CCD4E0
    "white_50":     RGBColor(0x88, 0x96, 0xAA),   # mid-grey text    #8896AA
    "graphite":     RGBColor(0x1E, 0x2D, 0x4A),   # graphite card    #1E2D4A
    "accent_line":  RGBColor(0x00, 0xD4, 0xFF),   # cyan rule
    "success":      RGBColor(0x00, 0xE5, 0xA0),   # teal accent      #00E5A0
    "warn":         RGBColor(0xFF, 0xB8, 0x20),   # amber            #FFB820
}

# Font — will use Calibri (universal) with size hierarchy
FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"
FONT_MONO    = "Consolas"


# ---------------------------------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor, left=0, top=0, width=None, height=None):
    """Fill slide background (or a rect) with a solid color."""
    w = width  or W
    h = height or H
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(left), Emu(top), Emu(w), Emu(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def rect(slide, left, top, width, height, fill: RGBColor, alpha=None,
         line_color: RGBColor = None, line_width_pt=0.75):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def line_h(slide, left, top, width, color: RGBColor, weight_pt=1.0):
    """Horizontal line."""
    connector = slide.shapes.add_connector(1, Emu(left), Emu(top),
                                            Emu(left + width), Emu(top))
    connector.line.color.rgb = color
    connector.line.width = Pt(weight_pt)
    return connector


def line_v(slide, left, top, height, color: RGBColor, weight_pt=1.0):
    connector = slide.shapes.add_connector(1, Emu(left), Emu(top),
                                            Emu(left), Emu(top + height))
    connector.line.color.rgb = color
    connector.line.width = Pt(weight_pt)
    return connector


def txb(slide, left, top, width, height, text: str, size_pt: float,
        color: RGBColor = None, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font=FONT_TITLE, wrap=True,
        word_wrap=True, space_after=0, space_before=0) -> object:
    """Add a text box with a single paragraph."""
    col = color or C["white"]
    tf = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf.text_frame.word_wrap = word_wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size_pt)
    run.font.color.rgb = col
    run.font.bold  = bold
    run.font.italic = italic
    return tf


def txb_ml(slide, left, top, width, height, lines: list[dict],
           word_wrap=True) -> object:
    """Multi-line text box. lines = [{text, size, color, bold, italic, align, space_after, space_before}]"""
    tf_shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = tf_shape.text_frame
    tf.word_wrap = word_wrap
    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_after  = Pt(ln.get("space_after",  0))
        p.space_before = Pt(ln.get("space_before", 0))
        run = p.add_run()
        run.text = ln.get("text", "")
        run.font.name  = ln.get("font", FONT_TITLE)
        run.font.size  = Pt(ln.get("size", 14))
        run.font.color.rgb = ln.get("color", C["white"])
        run.font.bold   = ln.get("bold",   False)
        run.font.italic = ln.get("italic", False)
    return tf_shape


def slide_label(slide, number: str, label: str):
    """Top-right slide label like  01  PROBLEM STATEMENT"""
    txb(slide, Inches(10.8), Inches(0.22), Inches(2.2), Inches(0.4),
        f"{number}  {label}", 7.5, color=C["white_50"],
        bold=False, align=PP_ALIGN.RIGHT, font=FONT_BODY)


def section_rule(slide, y_inch=0.9):
    """Thin cyan horizontal rule below the header area."""
    line_h(slide, Inches(0.55), Inches(y_inch), Inches(12.23), C["cyan_dim"], 0.5)


def slide_footer(slide, text="HYMIND   |   Hydrogen Market Intelligence & Data",
                 page_num=""):
    """Footer bar at the bottom of every slide."""
    rect(slide, 0, int(Inches(7.18)), int(W), int(Inches(0.32)), C["bg_card"])
    txb(slide, Inches(0.55), Inches(7.2), Inches(9), Inches(0.28),
        text, 6.5, color=C["white_50"], font=FONT_BODY)
    if page_num:
        txb(slide, Inches(12.3), Inches(7.2), Inches(0.8), Inches(0.28),
            page_num, 6.5, color=C["white_50"], align=PP_ALIGN.RIGHT, font=FONT_BODY)


def cyan_bullet(slide, x_in, y_in, text, size_pt=12.5, color=None, indent=0):
    """Single bullet with cyan dot."""
    col = color or C["white_80"]
    # bullet dot
    rect(slide,
         int(Inches(x_in + indent)), int(Inches(y_in + 0.06)),
         int(Inches(0.07)), int(Inches(0.07)), C["cyan"])
    txb(slide,
        int(Inches(x_in + indent + 0.15)), int(Inches(y_in)),
        int(Inches(5.5 - indent)), int(Inches(0.32)),
        text, size_pt, color=col, font=FONT_BODY)


def card(slide, left_in, top_in, w_in, h_in,
         title="", body_lines=None, title_color=None,
         fill=None, border_color=None):
    """Rounded-look info card (rectangle + title + body lines)."""
    fill_c  = fill         or C["graphite"]
    border  = border_color or C["blue_dim"]
    tc      = title_color  or C["cyan"]

    rect(slide, int(Inches(left_in)), int(Inches(top_in)),
         int(Inches(w_in)),   int(Inches(h_in)),
         fill_c, line_color=border, line_width_pt=0.6)

    if title:
        txb(slide, int(Inches(left_in + 0.18)), int(Inches(top_in + 0.14)),
            int(Inches(w_in - 0.3)), int(Inches(0.38)),
            title, 11, color=tc, bold=True, font=FONT_TITLE)

    if body_lines:
        y = top_in + (0.56 if title else 0.16)
        for bl in body_lines:
            txb(slide,
                int(Inches(left_in + 0.18)), int(Inches(y)),
                int(Inches(w_in - 0.3)),   int(Inches(0.3)),
                bl, 9.5, color=C["white_80"], font=FONT_BODY)
            y += 0.26

    return


def arrow_right(slide, x_in, y_in, length_in=0.4, color=None):
    """Simple right-pointing arrow line."""
    col = color or C["cyan"]
    line_h(slide, int(Inches(x_in)), int(Inches(y_in)),
           int(Inches(length_in)), col, 1.5)
    # arrowhead triangle (approximated with small rect)
    rect(slide,
         int(Inches(x_in + length_in - 0.04)),
         int(Inches(y_in - 0.055)),
         int(Inches(0.08)), int(Inches(0.11)), col)


def pipe_node(slide, x_in, y_in, w_in, h_in, label, sublabel="",
              fill=None, label_size=10):
    """Box for pipeline node."""
    fc = fill or C["bg_mid"]
    rect(slide, int(Inches(x_in)), int(Inches(y_in)),
         int(Inches(w_in)), int(Inches(h_in)),
         fc, line_color=C["cyan_dim"], line_width_pt=0.8)
    cy = y_in + (h_in - 0.22) / 2 if not sublabel else y_in + 0.06
    txb(slide, int(Inches(x_in + 0.08)), int(Inches(cy)),
        int(Inches(w_in - 0.16)), int(Inches(0.28)),
        label, label_size, color=C["white"], bold=True,
        font=FONT_TITLE, align=PP_ALIGN.CENTER)
    if sublabel:
        txb(slide, int(Inches(x_in + 0.08)),
            int(Inches(y_in + h_in - 0.32)),
            int(Inches(w_in - 0.16)), int(Inches(0.26)),
            sublabel, 7.5, color=C["white_50"],
            font=FONT_BODY, align=PP_ALIGN.CENTER)


def tech_badge(slide, x_in, y_in, label, sublabel="", fill=None):
    """Technology badge card."""
    fc = fill or C["bg_mid"]
    rect(slide, int(Inches(x_in)), int(Inches(y_in)),
         int(Inches(1.55)), int(Inches(1.05)),
         fc, line_color=C["blue_dim"], line_width_pt=0.6)
    # top accent line
    rect(slide, int(Inches(x_in)), int(Inches(y_in)),
         int(Inches(1.55)), int(Inches(0.04)), C["cyan"])
    txb(slide, int(Inches(x_in + 0.1)), int(Inches(y_in + 0.16)),
        int(Inches(1.35)), int(Inches(0.36)),
        label, 12.5, color=C["white"], bold=True, font=FONT_TITLE,
        align=PP_ALIGN.CENTER)
    if sublabel:
        txb(slide, int(Inches(x_in + 0.08)), int(Inches(y_in + 0.6)),
            int(Inches(1.39)), int(Inches(0.34)),
            sublabel, 7.5, color=C["white_50"], font=FONT_BODY,
            align=PP_ALIGN.CENTER)


def phase_badge(slide, x_in, y_in, num, label, items):
    """Agile phase badge."""
    rect(slide, int(Inches(x_in)), int(Inches(y_in)),
         int(Inches(1.72)), int(Inches(2.8)),
         C["bg_card"], line_color=C["blue_dim"], line_width_pt=0.6)
    # Phase number
    rect(slide, int(Inches(x_in)), int(Inches(y_in)),
         int(Inches(1.72)), int(Inches(0.48)), C["blue_dim"])
    txb(slide, int(Inches(x_in)), int(Inches(y_in + 0.04)),
        int(Inches(1.72)), int(Inches(0.38)),
        num, 17, color=C["cyan"], bold=True,
        align=PP_ALIGN.CENTER, font=FONT_TITLE)
    txb(slide, int(Inches(x_in + 0.08)), int(Inches(y_in + 0.54)),
        int(Inches(1.56)), int(Inches(0.34)),
        label, 9, color=C["white"], bold=True,
        font=FONT_TITLE, align=PP_ALIGN.CENTER)
    y = y_in + 1.0
    for it in items:
        cyan_bullet(slide, x_in + 0.12, y, it, size_pt=8, color=C["white_80"])
        y += 0.27


# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------

def slide_01_cover(prs):
    """Cover — HYMIND title, subtitle, premium dark hydrogen atmosphere."""
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])

    # Left dark-panel gradient suggestion (solid deep navy strip)
    rect(sl, 0, 0, int(Inches(5.4)), int(H), C["bg_section"])

    # Vertical cyan accent bar left edge
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])

    # Right: subtle mid-tone grid-feel overlay
    rect(sl, int(Inches(5.4)), 0, int(W) - int(Inches(5.4)), int(H), C["bg_deep"])

    # Top-right subtle box grid overlay (decorative)
    for i in range(6):
        for j in range(5):
            rx = Inches(5.7 + i * 1.22)
            ry = Inches(0.2 + j * 1.38)
            rect(sl, int(rx), int(ry), int(Inches(1.1)), int(Inches(1.22)),
                 C["bg_mid"])

    # Cyan rule on left panel
    line_h(sl, int(Inches(0.35)), int(Inches(4.82)), int(Inches(4.7)), C["cyan"], 1.5)

    # Brand prefix
    txb(sl, int(Inches(0.35)), int(Inches(1.6)), int(Inches(4.7)), int(Inches(0.5)),
        "HYDROGEN MARKET INTELLIGENCE & DATA",
        8.5, color=C["cyan"], bold=False, font=FONT_BODY,
        align=PP_ALIGN.LEFT)

    # Main title
    txb_ml(sl, int(Inches(0.35)), int(Inches(2.05)), int(Inches(4.8)),
           int(Inches(2.0)),
           [
               {"text": "HYMIND", "size": 54, "color": C["white"],
                "bold": True, "space_after": 0},
           ])

    # Divider between title and subtitle
    line_h(sl, int(Inches(0.35)), int(Inches(3.52)), int(Inches(2.2)), C["cyan_dim"], 0.75)

    # Subtitle
    txb(sl, int(Inches(0.35)), int(Inches(3.68)), int(Inches(4.8)), int(Inches(0.75)),
        "Autonomous Hydrogen Market\nIntelligence & Data Platform",
        14.5, color=C["white_80"], bold=False, font=FONT_BODY,
        align=PP_ALIGN.LEFT)

    # Key labels beneath rule
    y = 5.12
    for kw in ["Research Automation", "LangGraph Orchestration",
                "RAG Intelligence", "Executive Reporting"]:
        txb_ml(sl, int(Inches(0.35)), int(Inches(y)), int(Inches(4.8)),
               int(Inches(0.28)),
               [{"text": f"— {kw}", "size": 9.5, "color": C["white_50"],
                 "bold": False}])
        y += 0.30

    # Right side large HYMIND watermark text
    txb(sl, int(Inches(5.2)), int(Inches(1.6)), int(Inches(7.8)), int(Inches(3.5)),
        "HYMIND",
        96, color=RGBColor(0x0A, 0x18, 0x36), bold=True, font=FONT_TITLE,
        align=PP_ALIGN.CENTER)

    # Right side: tagline
    txb(sl, int(Inches(5.5)), int(Inches(5.0)), int(Inches(7.5)), int(Inches(0.45)),
        "Hydrogen Intelligence. Automated. Executive-Ready.",
        13, color=C["white_50"], bold=False, font=FONT_BODY,
        align=PP_ALIGN.CENTER)

    # Bottom footer
    slide_footer(sl, "CONFIDENTIAL — HYMIND Executive Presentation  2026")
    return sl


def slide_02_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "01", "PROBLEM STATEMENT")
    section_rule(sl, 1.08)

    # Header
    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "The Challenge", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(8.5)), int(Inches(0.42)),
        "Hydrogen industry stakeholders face fragmented, high-volume intelligence needs\n"
        "that are impossible to address manually at executive scale.",
        11.5, color=C["white_80"], font=FONT_BODY)

    # 5 problem cards in a row
    problems = [
        ("01", "Fragmented\nInformation",
         "Market signals scattered across news, feeds,\nwebsites, policy documents, and databases."),
        ("02", "Competitor\nBlind Spots",
         "No systematic way to track competitor\nannouncements, investments, or strategies."),
        ("03", "Policy &\nFunding Complexity",
         "Hydrogen regulation and government funding\nchange rapidly across multiple jurisdictions."),
        ("04", "Manual\nResearch Load",
         "Analyst time consumed by low-value data\ncollection rather than strategic interpretation."),
        ("05", "Delayed\nExecutive Insight",
         "Decision-makers receive weekly briefs days\nafter market-moving events have occurred."),
    ]

    x = 0.55
    for num, title, body in problems:
        rect(sl, int(Inches(x)), int(Inches(1.88)), int(Inches(2.38)),
             int(Inches(4.72)), C["bg_card"], line_color=C["blue_dim"], line_width_pt=0.6)
        # top accent
        rect(sl, int(Inches(x)), int(Inches(1.88)), int(Inches(2.38)),
             int(Inches(0.05)), C["cyan"])
        txb(sl, int(Inches(x + 0.15)), int(Inches(2.0)),
            int(Inches(2.1)), int(Inches(0.42)),
            num, 22, color=C["cyan"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(x + 0.15)), int(Inches(2.48)),
            int(Inches(2.1)), int(Inches(0.52)),
            title, 12, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(x + 0.15)), int(Inches(3.14)),
            int(Inches(2.1)), int(Inches(1.8)),
            body, 9.5, color=C["white_80"], font=FONT_BODY)
        x += 2.50

    slide_footer(sl, page_num="02")
    return sl


def slide_03_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "02", "SOLUTION OVERVIEW")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "HYMIND — The Solution", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(8.5)), int(Inches(0.38)),
        "An autonomous multi-source intelligence platform that collects, synthesizes, "
        "and delivers executive-grade hydrogen market intelligence on demand.",
        11.5, color=C["white_80"], font=FONT_BODY)

    # Left: what it does — large text
    txb_ml(sl, int(Inches(0.55)), int(Inches(1.72)),
           int(Inches(4.4)), int(Inches(4.9)),
           [
               {"text": "Autonomous\nMarket Intelligence",
                "size": 30, "color": C["white"], "bold": True, "space_after": 8},
               {"text": "for the Hydrogen Industry.", "size": 17,
                "color": C["cyan"], "bold": False, "space_after": 14},
               {"text": "HYMIND operates as a fully autonomous research pipeline:\n"
                        "collecting signals, validating sources, retrieving historical\n"
                        "context, and synthesising structured executive reports —\n"
                        "distributed automatically to decision-makers.",
                "size": 11, "color": C["white_80"], "bold": False},
           ])

    # Right: 4 capability cards
    caps = [
        ("Collect", "Serper · NewsAPI · RSS · Web Crawling"),
        ("Synthesise", "LangGraph · OpenAI GPT-4o-mini"),
        ("Remember", "Pinecone RAG · Historical Context"),
        ("Distribute", "FastAPI · n8n · Gmail · Sheets"),
    ]
    x0, y0 = 5.5, 1.72
    for i, (title, sub) in enumerate(caps):
        xi = x0 + (i % 2) * 3.65
        yi = y0 + (i // 2) * 1.62
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(3.38)),
             int(Inches(1.42)), C["graphite"],
             line_color=C["blue_dim"], line_width_pt=0.6)
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(3.38)),
             int(Inches(0.05)), C["cyan"])
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.14)),
            int(Inches(3.0)), int(Inches(0.38)),
            title, 14, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.56)),
            int(Inches(3.0)), int(Inches(0.56)),
            sub, 10, color=C["white_80"], font=FONT_BODY)

    slide_footer(sl, page_num="03")
    return sl


def slide_04_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "03", "ARCHITECTURE")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "System Architecture", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # Architecture pipeline — horizontal flow
    # 7 nodes across the slide
    nodes = [
        ("Trigger", "Manual / Scheduled\n/ HTTP API"),
        ("Collect", "Serper · News\nRSS · Crawl"),
        ("Validate", "Schema · Dedup\n· Normalise"),
        ("RAG\nStore", "Pinecone\nEmbeddings"),
        ("RAG\nRetrieve", "Historical\nContext"),
        ("Synthesise", "OpenAI\nGPT-4o-mini"),
        ("Distribute", "n8n · Gmail\n· Sheets"),
    ]

    n = len(nodes)
    box_w = 1.52
    gap = 0.22
    total = n * box_w + (n - 1) * gap
    x0 = (13.33 - total) / 2

    y_top = 1.55
    box_h = 1.5

    # Background band
    rect(sl, int(Inches(0.4)), int(Inches(y_top - 0.15)),
         int(Inches(12.53)), int(Inches(box_h + 0.55)),
         C["bg_section"])

    for idx, (label, sub) in enumerate(nodes):
        xi = x0 + idx * (box_w + gap)
        fill_c = C["blue_dim"] if idx == 0 else (C["bg_mid"] if idx < 6 else C["blue_dim"])
        pipe_node(sl, xi, y_top, box_w, box_h, label, sub, fill=fill_c, label_size=11)
        if idx < n - 1:
            arrow_right(sl, xi + box_w + 0.02, y_top + box_h / 2 - 0.02, gap - 0.04)

    # Three layer labels below
    layers = [
        (x0, "Input Layer", 1.52),
        (x0 + (box_w + gap), "Collection & Validation", 3 * box_w + 2 * gap),
        (x0 + 4 * (box_w + gap), "RAG Memory Layer", 2 * box_w + gap),
        (x0 + 6 * (box_w + gap), "Distribution", 1.52),
    ]
    for lx, lbl, lw in layers:
        txb(sl, int(Inches(lx)), int(Inches(3.28)),
            int(Inches(lw)), int(Inches(0.3)),
            lbl, 8.5, color=C["white_50"],
            align=PP_ALIGN.CENTER, font=FONT_BODY)
        line_h(sl, int(Inches(lx)), int(Inches(3.6)),
               int(Inches(lw)), C["blue_dim"], 0.5)

    # Below: 3 key architecture principles
    principles = [
        ("Modular Design",
         "Each pipeline layer is independently\ntestable and replaceable."),
        ("Graceful Degradation",
         "Missing APIs or credentials produce\nwarnings, not crashes."),
        ("State-Driven Orchestration",
         "LangGraph maintains typed state\nacross all 9 workflow nodes."),
    ]
    x = 0.55
    for title, body in principles:
        rect(sl, int(Inches(x)), int(Inches(3.98)), int(Inches(3.92)),
             int(Inches(2.68)), C["bg_card"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        txb(sl, int(Inches(x + 0.2)), int(Inches(4.12)),
            int(Inches(3.5)), int(Inches(0.36)),
            title, 12, color=C["cyan"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(x + 0.2)), int(Inches(4.55)),
            int(Inches(3.5)), int(Inches(1.8)),
            body, 10, color=C["white_80"], font=FONT_BODY)
        x += 4.16

    slide_footer(sl, page_num="04")
    return sl


def slide_05_pipeline(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "04", "RESEARCH PIPELINE")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Multi-Source Research Pipeline", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # Left: source icons / cards
    sources = [
        ("Serper API", "Google-based web research\nOrganic + news results"),
        ("NewsAPI", "Structured news retrieval\nFiltered, rate-limited"),
        ("RSS Feeds", "Industry feeds — H2 Insight,\nFuel Cells Works, H2 View"),
        ("Web Crawler", "Full-page content extraction\nrequests + BeautifulSoup"),
    ]
    y = 1.56
    for src, desc in sources:
        rect(sl, int(Inches(0.55)), int(Inches(y)), int(Inches(2.9)),
             int(Inches(1.02)), C["graphite"],
             line_color=C["cyan_dim"], line_width_pt=0.5)
        txb(sl, int(Inches(0.74)), int(Inches(y + 0.1)),
            int(Inches(2.5)), int(Inches(0.34)),
            src, 11.5, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(0.74)), int(Inches(y + 0.48)),
            int(Inches(2.5)), int(Inches(0.42)),
            desc, 8.5, color=C["white_80"], font=FONT_BODY)
        y += 1.18

    # Arrow pointing right from sources to merge
    line_h(sl, int(Inches(3.52)), int(Inches(3.6)),
           int(Inches(0.6)), C["cyan"], 1.5)
    rect(sl, int(Inches(4.05)), int(Inches(3.52)),
         int(Inches(0.08)), int(Inches(0.15)), C["cyan"])

    # Central processing column
    procs = [
        ("Merge & Deduplicate", "Normalised URL comparison\nURL-lowercase + trailing-slash stripped"),
        ("Schema Validation", "CollectorProtocol enforcement\nMissing URL / duplicate removal"),
        ("Crawl Selected", "Top 5 non-PDF URLs\nBoilerplate stripped"),
    ]
    y = 1.56
    for proc, desc in procs:
        rect(sl, int(Inches(4.2)), int(Inches(y)), int(Inches(3.6)),
             int(Inches(1.52)), C["bg_mid"],
             line_color=C["blue_dim"], line_width_pt=0.7)
        rect(sl, int(Inches(4.2)), int(Inches(y)), int(Inches(3.6)),
             int(Inches(0.05)), C["blue"])
        txb(sl, int(Inches(4.4)), int(Inches(y + 0.1)),
            int(Inches(3.2)), int(Inches(0.38)),
            proc, 11.5, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(4.4)), int(Inches(y + 0.54)),
            int(Inches(3.2)), int(Inches(0.72)),
            desc, 9.5, color=C["white_80"], font=FONT_BODY)
        if y < 4.0:
            line_v(sl, int(Inches(6.0)), int(Inches(y + 1.52)),
                   int(Inches(0.3)), C["cyan_dim"], 0.8)
        y += 1.82

    # Arrow from processing to output
    line_h(sl, int(Inches(7.88)), int(Inches(3.6)),
           int(Inches(0.6)), C["cyan"], 1.5)
    rect(sl, int(Inches(8.41)), int(Inches(3.52)),
         int(Inches(0.08)), int(Inches(0.15)), C["cyan"])

    # Right: output schema
    rect(sl, int(Inches(8.54)), int(Inches(1.56)), int(Inches(4.24)),
         int(Inches(5.22)), C["bg_card"],
         line_color=C["blue_dim"], line_width_pt=0.6)
    rect(sl, int(Inches(8.54)), int(Inches(1.56)), int(Inches(4.24)),
         int(Inches(0.46)), C["blue_dim"])
    txb(sl, int(Inches(8.74)), int(Inches(1.62)),
        int(Inches(3.84)), int(Inches(0.36)),
        "Unified Output Schema", 11, color=C["white"],
        bold=True, font=FONT_TITLE)

    fields = [
        ("title", "Article / result title"),
        ("url", "Source URL"),
        ("snippet", "Summary / description"),
        ("published_at", "Publication date"),
        ("source", "Publisher name"),
        ("source_type", "organic · news · rss"),
        ("search_query", "Query or feed URL"),
        ("rank", "Position within source"),
    ]
    y = 2.16
    for field, desc in fields:
        txb(sl, int(Inches(8.74)), int(Inches(y)),
            int(Inches(1.3)), int(Inches(0.26)),
            field, 9, color=C["cyan"], bold=True, font=FONT_MONO)
        txb(sl, int(Inches(10.14)), int(Inches(y)),
            int(Inches(2.44)), int(Inches(0.26)),
            desc, 9, color=C["white_80"], font=FONT_BODY)
        y += 0.32

    slide_footer(sl, page_num="05")
    return sl


def slide_06_rag(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "05", "RAG & INTELLIGENCE LAYER")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "RAG — Persistent Market Memory", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(9.5)), int(Inches(0.38)),
        "Every research run enriches a Pinecone vector database, enabling historical trend "
        "comparison, contextual retrieval, and persistent intelligence across weekly cycles.",
        11, color=C["white_80"], font=FONT_BODY)

    # RAG flow diagram
    rag_steps = [
        ("Findings\nMerged", "Post-crawl\nresults"),
        ("Embed", "text-embedding\n-3-small"),
        ("Upsert", "Pinecone\nServerless"),
        ("Retrieve", "Top-5 semantic\nmatches"),
        ("Inject", "Historical\nContext"),
        ("Synthesise", "OpenAI\nReport"),
    ]
    box_w, box_h = 1.72, 1.2
    gap = 0.22
    x0 = 0.55
    y0 = 1.76

    for idx, (lbl, sub) in enumerate(rag_steps):
        fill_c = C["bg_mid"] if idx not in (2, 3) else C["blue_dim"]
        pipe_node(sl, x0 + idx * (box_w + gap), y0, box_w, box_h, lbl, sub,
                  fill=fill_c, label_size=10.5)
        if idx < len(rag_steps) - 1:
            arrow_right(sl,
                        x0 + idx * (box_w + gap) + box_w + 0.02,
                        y0 + box_h / 2 - 0.02,
                        gap - 0.04)

    # Pinecone annotation
    txb(sl, int(Inches(3.72)), int(Inches(3.12)), int(Inches(3.64)), int(Inches(0.3)),
        "↑  Pinecone Index  |  1536-dim cosine  |  AWS us-east-1",
        8.5, color=C["cyan_dim"], font=FONT_BODY, align=PP_ALIGN.CENTER)

    # 4 benefit cards
    benefits = [
        ("Historical Awareness",
         "Every run builds on previous intelligence cycles. "
         "Trend signals strengthen over time."),
        ("Source Traceability",
         "Each stored finding preserves url, source, published_at, "
         "topic, and collected_at metadata."),
        ("Semantic Retrieval",
         "Top-5 cosine-similar historical findings are injected into "
         "the synthesis prompt automatically."),
        ("Graceful Degradation",
         "Pipeline runs normally without Pinecone configured — "
         "RAG nodes emit a warning and continue."),
    ]
    x = 0.55
    for title, body in benefits:
        rect(sl, int(Inches(x)), int(Inches(3.58)), int(Inches(2.96)),
             int(Inches(2.88)), C["bg_card"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        rect(sl, int(Inches(x)), int(Inches(3.58)), int(Inches(2.96)),
             int(Inches(0.05)), C["cyan"])
        txb(sl, int(Inches(x + 0.18)), int(Inches(3.7)),
            int(Inches(2.6)), int(Inches(0.38)),
            title, 11, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(x + 0.18)), int(Inches(4.14)),
            int(Inches(2.6)), int(Inches(2.0)),
            body, 9.5, color=C["white_80"], font=FONT_BODY)
        x += 3.15

    slide_footer(sl, page_num="06")
    return sl


def slide_07_langgraph(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "06", "LANGGRAPH WORKFLOW")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "LangGraph Agent Orchestration", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # Left: 9-node pipeline vertical list
    nodes = [
        ("initialize_state",         "Sets run metadata, timestamps, topic"),
        ("collect_serper",           "Serper organic + news results"),
        ("collect_news",             "NewsAPI articles with retry"),
        ("collect_rss",              "Hydrogen industry RSS/Atom feeds"),
        ("merge_and_deduplicate",    "URL-normalised deduplication"),
        ("crawl_selected",           "Full content extraction, top 5 URLs"),
        ("store_findings_in_pinecone", "Embed + upsert to Pinecone"),
        ("retrieve_context_from_pinecone", "Top-5 semantic history"),
        ("finalize_state",           "Counts, duration, error summary"),
    ]

    y = 1.52
    for i, (node, desc) in enumerate(nodes):
        fc = C["blue_dim"] if i in (0, 8) else C["bg_mid"]
        rect(sl, int(Inches(0.55)), int(Inches(y)), int(Inches(3.5)),
             int(Inches(0.54)), fc, line_color=C["cyan_dim"], line_width_pt=0.5)
        txb(sl, int(Inches(0.72)), int(Inches(y + 0.06)),
            int(Inches(3.1)), int(Inches(0.28)),
            node, 9, color=C["cyan"] if i in (0, 8) else C["white"],
            bold=True, font=FONT_MONO)
        txb(sl, int(Inches(0.72)), int(Inches(y + 0.3)),
            int(Inches(3.1)), int(Inches(0.2)),
            desc, 7.5, color=C["white_50"], font=FONT_BODY)
        if i < 8:
            line_v(sl, int(Inches(0.92)), int(Inches(y + 0.54)),
                   int(Inches(0.14)), C["cyan_dim"], 0.8)
        y += 0.68

    # Right: 4 LangGraph capability cards
    rect(sl, int(Inches(4.6)), int(Inches(1.52)), int(Inches(8.3)),
         int(Inches(5.38)), C["bg_section"])

    lg_caps = [
        ("Typed State Machine",
         "AgentState TypedDict flows through all 9 nodes.\n"
         "errors and warnings use Annotated + operator.add\n"
         "for safe concurrent accumulation."),
        ("Per-Node Isolation",
         "Every collection node catches ALL exceptions\n"
         "and appends to state[\"errors\"]. One failing\n"
         "source never blocks the others."),
        ("=== Node START / END === Logging",
         "Each node emits structured START and END log\n"
         "markers with counts and duration. Unambiguous\n"
         "pipeline visibility during demo and production."),
        ("Graceful Continuation",
         "Missing API keys are converted to warnings\n"
         "(not sys.exit). The pipeline always completes\n"
         "with whatever partial results are available."),
    ]
    x0, y0 = 4.82, 1.64
    for i, (title, body) in enumerate(lg_caps):
        xi = x0 + (i % 2) * 4.0
        yi = y0 + (i // 2) * 2.56
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(3.72)),
             int(Inches(2.34)), C["graphite"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(3.72)),
             int(Inches(0.05)), C["cyan"])
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.12)),
            int(Inches(3.36)), int(Inches(0.38)),
            title, 11, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.54)),
            int(Inches(3.36)), int(Inches(1.6)),
            body, 9.5, color=C["white_80"], font=FONT_BODY)

    slide_footer(sl, page_num="07")
    return sl


def slide_08_reliability(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "07", "RELIABILITY & HARDENING")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Production-Grade Reliability", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(9.5)), int(Inches(0.38)),
        "Phase 4 hardened every external integration and internal pipeline stage against "
        "real-world failure conditions.",
        11, color=C["white_80"], font=FONT_BODY)

    # Left column: test results
    rect(sl, int(Inches(0.55)), int(Inches(1.72)), int(Inches(4.0)),
         int(Inches(5.08)), C["bg_card"],
         line_color=C["blue_dim"], line_width_pt=0.6)
    rect(sl, int(Inches(0.55)), int(Inches(1.72)), int(Inches(4.0)),
         int(Inches(0.48)), C["blue_dim"])
    txb(sl, int(Inches(0.75)), int(Inches(1.78)), int(Inches(3.6)),
        int(Inches(0.38)), "Test Coverage", 12,
        color=C["white"], bold=True, font=FONT_TITLE)

    test_items = [
        ("243", "Tests — all pass, no live API calls"),
        ("73",  "Failure scenario tests (Phase 4)"),
        ("46",  "RAG layer tests"),
        ("60",  "NewsAPI collector tests"),
        ("34",  "Validator unit tests"),
        ("20",  "Web crawler tests"),
    ]
    y = 2.36
    for num, desc in test_items:
        txb(sl, int(Inches(0.75)), int(Inches(y)),
            int(Inches(0.72)), int(Inches(0.36)),
            num, 18, color=C["cyan"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(1.54)), int(Inches(y + 0.06)),
            int(Inches(2.8)), int(Inches(0.3)),
            desc, 9.5, color=C["white_80"], font=FONT_BODY)
        y += 0.62

    # Right: reliability feature cards 2x3
    features = [
        ("Tenacity Retry",
         "Exponential back-off on Timeout, ConnectionError, HTTP 429, HTTP 5xx for every external call."),
        ("Rate Limit Protection",
         "NewsAPI and Serper failures handled with retry — pipeline continues with partial results."),
        ("Schema Validation",
         "CollectorProtocol enforces output schema. Missing-URL and duplicate entries removed pre-LLM."),
        ("PDF Filtering",
         "PDF URLs automatically excluded from crawl queue. Boilerplate stripped pre-extraction."),
        ("No Secrets in Logs",
         "API keys sent via HTTP headers, never in URL params. Keys are never written to log files."),
        ("Node-Level Logging",
         "=== Node START/END === markers with counts on every node. Clear pipeline observability."),
    ]
    x0, y0 = 4.82, 1.72
    for i, (title, body) in enumerate(features):
        xi = x0 + (i % 2) * 4.06
        yi = y0 + (i // 2) * 1.68
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(3.8)),
             int(Inches(1.5)), C["graphite"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.1)),
            int(Inches(3.44)), int(Inches(0.34)),
            title, 11, color=C["cyan"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(xi + 0.18)), int(Inches(yi + 0.48)),
            int(Inches(3.44)), int(Inches(0.88)),
            body, 9.5, color=C["white_80"], font=FONT_BODY)

    slide_footer(sl, page_num="08")
    return sl


def slide_09_n8n(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "08", "DISTRIBUTION AUTOMATION")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "n8n Distribution Workflow", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(9.5)), int(Inches(0.38)),
        "n8n orchestrates the weekly delivery cycle — from schedule trigger through "
        "Markdown-to-HTML conversion, Gmail delivery, and Google Sheets logging.",
        11, color=C["white_80"], font=FONT_BODY)

    # n8n flow — vertical steps on left
    n8n_steps = [
        ("Schedule Trigger",   "Every Monday 08:00"),
        ("HTTP Request",       "POST /run-hymind via ngrok"),
        ("IF — status check",  "$json.status == success"),
        ("Markdown → HTML",    "n8n built-in Markdown node"),
        ("Gmail — Send",       "HTML report to recipient"),
        ("Google Sheets",      "Delivery log: timestamp · title · status"),
    ]

    y = 1.72
    for i, (step, sub) in enumerate(n8n_steps):
        fc = C["bg_mid"] if i not in (0,) else C["blue_dim"]
        ic = C["cyan"] if i == 0 else (C["success"] if i == 4 else C["white"])
        rect(sl, int(Inches(0.55)), int(Inches(y)), int(Inches(3.9)),
             int(Inches(0.68)), fc, line_color=C["cyan_dim"], line_width_pt=0.5)
        # step number
        rect(sl, int(Inches(0.55)), int(Inches(y)), int(Inches(0.38)),
             int(Inches(0.68)), C["blue_dim"])
        txb(sl, int(Inches(0.55)), int(Inches(y + 0.18)),
            int(Inches(0.38)), int(Inches(0.3)),
            str(i + 1), 12, color=C["cyan"], bold=True,
            font=FONT_TITLE, align=PP_ALIGN.CENTER)
        txb(sl, int(Inches(1.02)), int(Inches(y + 0.04)),
            int(Inches(3.24)), int(Inches(0.3)),
            step, 11, color=ic, bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(1.02)), int(Inches(y + 0.34)),
            int(Inches(3.24)), int(Inches(0.28)),
            sub, 9, color=C["white_50"], font=FONT_BODY)
        if i < 5:
            line_v(sl, int(Inches(0.745)), int(Inches(y + 0.68)),
                   int(Inches(0.12)), C["cyan_dim"], 0.8)
        y += 0.82

    # Error branch annotation
    txb(sl, int(Inches(0.55)), int(Inches(6.38)), int(Inches(3.9)), int(Inches(0.38)),
        "Error branch → Gmail alert + Global Error Handler workflow",
        8.5, color=C["warn"], font=FONT_BODY, italic=True)

    # Right: architecture detail
    rect(sl, int(Inches(5.1)), int(Inches(1.52)), int(Inches(7.78)),
         int(Inches(5.28)), C["bg_section"])

    details = [
        ("FastAPI Wrapper",
         "src/api/server.py exposes POST /run-hymind and GET /health.\n"
         "Failures return HTTP 200 with status: failed so n8n reads the\n"
         "error without triggering its own HTTP error handling."),
        ("ngrok Tunnel",
         "ngrok creates a public HTTPS tunnel to localhost:8000.\n"
         "Fixed ngrok domain supported for stable URLs.\n"
         "start_hymind_api.py starts FastAPI + ngrok in one command."),
        ("PDF Decision",
         "PDF generation was evaluated and descoped from MVP.\n"
         "n8n's built-in Markdown node converts to HTML inline —\n"
         "email-readable format with zero additional dependencies."),
        ("Google Sheets Logging",
         "Every delivery writes: timestamp · report_title · status · channel.\n"
         "Delivery audit trail stored in Google Sheets via OAuth2.\n"
         "n8n/HYMIND.json is exported and committed to the repository."),
    ]
    y = 1.66
    for title, body in details:
        rect(sl, int(Inches(5.3)), int(Inches(y)), int(Inches(7.38)),
             int(Inches(1.18)), C["graphite"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        txb(sl, int(Inches(5.5)), int(Inches(y + 0.08)),
            int(Inches(7.0)), int(Inches(0.34)),
            title, 11, color=C["cyan"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(5.5)), int(Inches(y + 0.44)),
            int(Inches(7.0)), int(Inches(0.68)),
            body, 9, color=C["white_80"], font=FONT_BODY)
        y += 1.3

    slide_footer(sl, page_num="09")
    return sl


def slide_10_sample_report(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "09", "SAMPLE EXECUTIVE REPORT")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Executive Intelligence Report — Output Format", 26,
        color=C["white"], bold=True, font=FONT_TITLE)

    # Simulated report mockup (paper-on-dark)
    rect(sl, int(Inches(0.55)), int(Inches(1.38)), int(Inches(6.6)),
         int(Inches(5.78)), RGBColor(0x06, 0x12, 0x28),
         line_color=C["blue_dim"], line_width_pt=0.8)

    # Report title bar
    rect(sl, int(Inches(0.55)), int(Inches(1.38)), int(Inches(6.6)),
         int(Inches(0.52)), C["bg_mid"])
    txb(sl, int(Inches(0.75)), int(Inches(1.44)), int(Inches(6.2)),
        int(Inches(0.38)),
        "HYMIND Executive Intelligence Report",
        11, color=C["white"], bold=True, font=FONT_TITLE)

    # Report content mockup
    sections = [
        ("Research Topic",
         "European green hydrogen electrolyzer market capacity & cost trajectory 2026"),
        ("Executive Summary",
         "European electrolyzer deployment accelerated in H1 2026, driven by IRA-competitive "
         "policy responses from Brussels and large-scale FIDs. Installed capacity reached "
         "an estimated 1.8 GW across the EU. Production costs remain above 5–6 EUR/kg..."),
        ("Key Developments",
         "— German H2Global mechanism awarded six import contracts totalling 420,000 t/yr\n"
         "— Nel Hydrogen announced second 1 GW/year gigafactory in Norway (Q4 2027)\n"
         "— Spain increased PERTE H2 budget by €800M for 12 corridor projects"),
        ("Policy & Funding",
         "EU REPowerEU 10 Mt domestic production target · US IRA competitive pressure\n"
         "EC revised Hydrogen Delegated Act · Offtake pricing 5–7 EUR/kg for 2027 delivery"),
    ]
    y = 2.04
    for title, content in sections:
        txb(sl, int(Inches(0.75)), int(Inches(y)), int(Inches(6.2)),
            int(Inches(0.28)), title,
            8, color=C["cyan"], bold=True, font=FONT_TITLE)
        line_h(sl, int(Inches(0.75)), int(Inches(y + 0.28)),
               int(Inches(6.2)), C["blue_dim"], 0.4)
        txb(sl, int(Inches(0.75)), int(Inches(y + 0.32)),
            int(Inches(6.2)), int(Inches(0.7)),
            content, 8, color=C["white_80"], font=FONT_BODY)
        y += 1.12

    txb(sl, int(Inches(0.75)), int(Inches(6.68)), int(Inches(6.2)),
        int(Inches(0.24)),
        "Generated by HYMIND — Hydrogen Market Intelligence & Data",
        7.5, color=C["white_50"], font=FONT_BODY, italic=True)

    # Right: report structure overview
    txb(sl, int(Inches(7.6)), int(Inches(1.38)), int(Inches(5.28)),
        int(Inches(0.38)),
        "Report Structure", 13, color=C["cyan"], bold=True, font=FONT_TITLE)
    line_h(sl, int(Inches(7.6)), int(Inches(1.78)), int(Inches(5.28)),
           C["cyan_dim"], 0.5)

    report_secs = [
        ("01", "Research Topic"),
        ("02", "Executive Summary"),
        ("03", "Key Developments"),
        ("04", "Market Implications"),
        ("05", "Technology Signals"),
        ("06", "Policy & Funding"),
        ("07", "Competitive Notes"),
        ("08", "Risks & Watchouts"),
        ("09", "Source Traceability"),
        ("10", "Workflow Metadata"),
    ]
    y = 1.94
    for num, title in report_secs:
        txb(sl, int(Inches(7.6)), int(Inches(y)), int(Inches(0.38)),
            int(Inches(0.32)), num, 9, color=C["cyan"],
            bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(8.1)), int(Inches(y)), int(Inches(4.6)),
            int(Inches(0.32)), title, 10.5, color=C["white_80"], font=FONT_BODY)
        y += 0.48

    slide_footer(sl, page_num="10")
    return sl


def slide_11_techstack(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "10", "TECHNICAL STACK")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Technology Stack", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # 8 tech badges — 2 rows of 4 (with more space for sub-labels)
    stack = [
        ("Python\n3.11+",   "Core Language",      C["bg_mid"]),
        ("LangGraph",        "Agent Orchestration", C["bg_mid"]),
        ("OpenAI",           "GPT-4o-mini / Embed", C["bg_mid"]),
        ("Pinecone",         "Vector RAG Memory",   C["blue_dim"]),
        ("n8n",              "Distribution Workflow",C["bg_mid"]),
        ("FastAPI",          "HTTP API Wrapper",    C["bg_mid"]),
        ("Serper + NewsAPI", "Research Sources",    C["bg_mid"]),
        ("feedparser +\nBeautifulSoup", "RSS + Crawling", C["bg_mid"]),
    ]

    cols = 4
    box_w, box_h = 2.88, 1.56
    gap_x, gap_y = 0.3, 0.42
    x0 = (13.33 - cols * box_w - (cols - 1) * gap_x) / 2
    y0 = 1.52

    for i, (label, sub, fill_c) in enumerate(stack):
        row = i // cols
        col = i % cols
        xi = x0 + col * (box_w + gap_x)
        yi = y0 + row * (box_h + gap_y)

        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(box_w)),
             int(Inches(box_h)), fill_c,
             line_color=C["blue_dim"], line_width_pt=0.6)
        # Top accent bar
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(box_w)),
             int(Inches(0.05)), C["cyan"])
        txb(sl, int(Inches(xi + 0.14)), int(Inches(yi + 0.18)),
            int(Inches(box_w - 0.28)), int(Inches(0.64)),
            label, 14, color=C["white"], bold=True, font=FONT_TITLE,
            align=PP_ALIGN.CENTER)
        txb(sl, int(Inches(xi + 0.14)), int(Inches(yi + 0.88)),
            int(Inches(box_w - 0.28)), int(Inches(0.48)),
            sub, 9.5, color=C["white_50"], font=FONT_BODY,
            align=PP_ALIGN.CENTER)

    slide_footer(sl, page_num="11")
    return sl


def slide_12_agile(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "11", "DEVELOPMENT APPROACH")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Agile Phased Development", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # Phase cards — 6 phases across
    phases = [
        ("0", "Foundation",    ["Repo scaffold", "Docs structure", "Memory system", "Skills setup"]),
        ("1", "Core Pipeline", ["OpenAI · Serper", "NewsAPI · RSS", "LangGraph workflow", "72 tests"]),
        ("2", "Validation",    ["CollectorProtocol", "Schema validation", "HTML stripping", "+60 tests"]),
        ("3", "RAG Memory",    ["Pinecone store", "Semantic retrieval", "History injection", "+46 tests"]),
        ("4", "Hardening",     ["Output validator", "73 failure tests", "Node logging", "243 total tests"]),
        ("5", "Distribution",  ["FastAPI wrapper", "n8n workflow", "Gmail delivery", "Sheets logging"]),
    ]

    box_w = 1.94
    gap   = 0.22
    x0    = (13.33 - 6 * box_w - 5 * gap) / 2
    y0    = 1.52

    for i, (num, label, items) in enumerate(phases):
        xi = x0 + i * (box_w + gap)
        rect(sl, int(Inches(xi)), int(Inches(y0)), int(Inches(box_w)),
             int(Inches(5.42)), C["bg_card"],
             line_color=C["blue_dim"], line_width_pt=0.6)
        # Phase header
        rect(sl, int(Inches(xi)), int(Inches(y0)), int(Inches(box_w)),
             int(Inches(0.58)), C["blue_dim"])
        txb(sl, int(Inches(xi)), int(Inches(y0 + 0.04)),
            int(Inches(box_w)), int(Inches(0.32)),
            f"Phase {num}", 16, color=C["cyan"], bold=True,
            font=FONT_TITLE, align=PP_ALIGN.CENTER)
        txb(sl, int(Inches(xi + 0.1)), int(Inches(y0 + 0.68)),
            int(Inches(box_w - 0.2)), int(Inches(0.34)),
            label, 10.5, color=C["white"], bold=True,
            font=FONT_TITLE, align=PP_ALIGN.CENTER)
        y = y0 + 1.14
        for it in items:
            # Bullet dot
            rect(sl, int(Inches(xi + 0.22)), int(Inches(y + 0.08)),
                 int(Inches(0.07)), int(Inches(0.07)), C["cyan"])
            txb(sl, int(Inches(xi + 0.36)), int(Inches(y)),
                int(Inches(box_w - 0.5)), int(Inches(0.3)),
                it, 9, color=C["white_80"], font=FONT_BODY)
            y += 0.4

        # Completion badge
        rect(sl, int(Inches(xi + 0.22)), int(Inches(y0 + 4.72)),
             int(Inches(box_w - 0.44)), int(Inches(0.38)),
             C["success"] if num != "6" else C["blue_dim"])
        txb(sl, int(Inches(xi + 0.22)), int(Inches(y0 + 4.76)),
            int(Inches(box_w - 0.44)), int(Inches(0.3)),
            "Complete" if num != "6" else "In Progress",
            8.5, color=C["bg_deep"], bold=True,
            font=FONT_BODY, align=PP_ALIGN.CENTER)

    slide_footer(sl, page_num="12")
    return sl


def slide_13_future(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "12", "FUTURE ROADMAP")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "Future Expansion", 26, color=C["white"], bold=True, font=FONT_TITLE)

    txb(sl, int(Inches(0.55)), int(Inches(1.18)), int(Inches(9.5)), int(Inches(0.38)),
        "The modular architecture positions HYMIND for enterprise-scale expansion "
        "across additional channels, industries, and intelligence dimensions.",
        11, color=C["white_80"], font=FONT_BODY)

    future_areas = [
        ("Enterprise\nIntegration",
         ["Microsoft Teams", "SharePoint connector", "Jira / PLM sync", "MS Graph API"]),
        ("Intelligence\nExpansion",
         ["LinkedIn monitoring", "Patent filings", "Conference signals", "Job market signals"]),
        ("Delivery\n& UX",
         ["Real-time dashboards", "Slack / Telegram", "PDF branded reports", "Web portal"]),
        ("Multi-Industry\nSupport",
         ["Battery / energy storage", "Clean energy broadly", "Automotive supply chain", "Electrolysis tech"]),
        ("AI Enhancements",
         ["Multi-model fallback", "Confidence scoring", "Human approval nodes", "Predictive analytics"]),
        ("Operations",
         ["CI/CD pipeline", "Prometheus metrics", "Grafana dashboards", "Load / chaos testing"]),
    ]

    x0, y0 = 0.55, 1.72
    box_w, box_h = 3.82, 2.42
    gap_x, gap_y = 0.25, 0.28

    for i, (title, items) in enumerate(future_areas):
        row = i // 3
        col = i % 3
        xi = x0 + col * (box_w + gap_x)
        yi = y0 + row * (box_h + gap_y)

        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(box_w)),
             int(Inches(box_h)), C["bg_card"],
             line_color=C["blue_dim"], line_width_pt=0.5)
        rect(sl, int(Inches(xi)), int(Inches(yi)), int(Inches(0.06)),
             int(Inches(box_h)), C["cyan_dim"])
        txb(sl, int(Inches(xi + 0.2)), int(Inches(yi + 0.14)),
            int(Inches(box_w - 0.3)), int(Inches(0.48)),
            title, 12, color=C["cyan"], bold=True, font=FONT_TITLE)
        y = yi + 0.7
        for it in items:
            rect(sl, int(Inches(xi + 0.2)), int(Inches(y + 0.08)),
                 int(Inches(0.06)), int(Inches(0.06)), C["white_50"])
            txb(sl, int(Inches(xi + 0.34)), int(Inches(y)),
                int(Inches(box_w - 0.5)), int(Inches(0.28)),
                it, 9.5, color=C["white_80"], font=FONT_BODY)
            y += 0.36

    slide_footer(sl, page_num="13")
    return sl


def slide_14_outcome(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_deep"])
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])
    slide_label(sl, "13", "PROJECT OUTCOME")
    section_rule(sl, 1.08)

    txb(sl, int(Inches(0.55)), int(Inches(0.22)), int(Inches(9)), int(Inches(0.55)),
        "What HYMIND Delivers", 26, color=C["white"], bold=True, font=FONT_TITLE)

    # Large central statement
    txb_ml(sl, int(Inches(0.55)), int(Inches(1.28)),
           int(Inches(12.23)), int(Inches(1.62)),
           [
               {"text": "A production-ready autonomous hydrogen market intelligence platform.",
                "size": 22, "color": C["white"], "bold": True,
                "align": PP_ALIGN.CENTER, "space_after": 6},
               {"text": "243 tests · 5 complete phases · FastAPI + LangGraph + Pinecone + n8n",
                "size": 12, "color": C["white_50"],
                "align": PP_ALIGN.CENTER},
           ])

    line_h(sl, int(Inches(2.0)), int(Inches(3.08)),
           int(Inches(9.33)), C["cyan_dim"], 0.6)

    # 4 outcome statements
    outcomes = [
        ("Autonomous\nOperation",
         "Runs end-to-end without human intervention — from collection through synthesis to delivery."),
        ("Executive-Grade\nReporting",
         "Structured 10-section Markdown reports with source traceability, RAG historical context, and metadata."),
        ("Enterprise-Ready\nArchitecture",
         "Modular LangGraph pipeline, typed state, graceful degradation, retry logic, 243-test coverage."),
        ("Scalable\nDelivery",
         "FastAPI HTTP wrapper + n8n enables weekly scheduled delivery, Gmail distribution, and audit logging."),
    ]

    x = 0.55
    for title, body in outcomes:
        rect(sl, int(Inches(x)), int(Inches(3.28)), int(Inches(2.9)),
             int(Inches(3.28)), C["bg_card"],
             line_color=C["blue_dim"], line_width_pt=0.6)
        rect(sl, int(Inches(x)), int(Inches(3.28)), int(Inches(2.9)),
             int(Inches(0.06)), C["cyan"])
        txb(sl, int(Inches(x + 0.2)), int(Inches(3.42)),
            int(Inches(2.5)), int(Inches(0.58)),
            title, 14, color=C["white"], bold=True, font=FONT_TITLE)
        txb(sl, int(Inches(x + 0.2)), int(Inches(4.06)),
            int(Inches(2.5)), int(Inches(2.2)),
            body, 10, color=C["white_80"], font=FONT_BODY)
        x += 3.2

    slide_footer(sl, page_num="14")
    return sl


def slide_15_closing(prs):
    sl = blank_slide(prs)
    fill_bg(sl, C["bg_section"])

    # Left cyan bar
    rect(sl, 0, 0, int(Inches(0.06)), int(H), C["cyan"])

    # Right dark panel
    rect(sl, int(Inches(0.06)), 0, int(W) - int(Inches(0.06)), int(H), C["bg_deep"])

    # Large background watermark
    txb(sl, int(Inches(0.3)), int(Inches(0.8)), int(Inches(12.7)), int(Inches(4.2)),
        "HYMIND",
        120, color=RGBColor(0x07, 0x14, 0x2A), bold=True, font=FONT_TITLE,
        align=PP_ALIGN.CENTER)

    # Cyan horizontal rule
    line_h(sl, int(Inches(3.2)), int(Inches(3.28)), int(Inches(6.93)),
           C["cyan"], 1.5)

    # HYMIND title
    txb(sl, int(Inches(0.55)), int(Inches(2.28)), int(Inches(12.23)),
        int(Inches(0.82)),
        "HYMIND", 52, color=C["white"], bold=True, font=FONT_TITLE,
        align=PP_ALIGN.CENTER)

    # Subtitle
    txb(sl, int(Inches(0.55)), int(Inches(3.42)), int(Inches(12.23)),
        int(Inches(0.5)),
        "Hydrogen Market Intelligence & Data",
        16, color=C["cyan"], bold=False, font=FONT_BODY,
        align=PP_ALIGN.CENTER)

    # Questions
    txb(sl, int(Inches(0.55)), int(Inches(4.22)), int(Inches(12.23)),
        int(Inches(0.44)),
        "Questions?",
        22, color=C["white_80"], bold=False, font=FONT_TITLE,
        align=PP_ALIGN.CENTER)

    # Bottom credit line
    txb(sl, int(Inches(0.55)), int(Inches(5.58)), int(Inches(12.23)),
        int(Inches(0.32)),
        "Autonomous Research · LangGraph · OpenAI · Pinecone · n8n · FastAPI",
        9, color=C["white_50"], font=FONT_BODY, align=PP_ALIGN.CENTER)

    # Footer
    slide_footer(sl, "CONFIDENTIAL — HYMIND Executive Presentation  2026")
    return sl


# ---------------------------------------------------------------------------
# BUILD
# ---------------------------------------------------------------------------

def build():
    prs = new_prs()

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_architecture(prs)
    slide_05_pipeline(prs)
    slide_06_rag(prs)
    slide_07_langgraph(prs)
    slide_08_reliability(prs)
    slide_09_n8n(prs)
    slide_10_sample_report(prs)
    slide_11_techstack(prs)
    slide_12_agile(prs)
    slide_13_future(prs)
    slide_14_outcome(prs)
    slide_15_closing(prs)

    out = Path("outputs/HYMIND_Executive_Presentation.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"[HYMIND] Presentation saved: {out}")
    return out


if __name__ == "__main__":
    build()
